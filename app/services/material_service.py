from pathlib import Path

from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation


def extract_text_from_material(file_path: str) -> str:
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    suffix = path.suffix.lower()

    if suffix == ".pdf":
        reader = PdfReader(str(path))
        text = []
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text.append(extracted)
        return "\n".join(text)

    elif suffix == ".docx":
        doc = Document(str(path))
        return "\n".join([para.text for para in doc.paragraphs])

    elif suffix == ".pptx":
        prs = Presentation(str(path))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    else:
        raise ValueError(f"Unsupported file type: {suffix}")
    

def save_extracted_text(file_path: str, text_dir: Path) -> str:
    extracted_text = extract_text_from_material(file_path).strip()

    if not extracted_text:
        raise ValueError("No text could be extracted from this file")

    original_path = Path(file_path)
    txt_filename = f"{original_path.stem}.txt"
    txt_path = text_dir / txt_filename

    text_dir.mkdir(parents=True, exist_ok=True)
    txt_path.write_text(extracted_text, encoding="utf-8")

    return str(txt_path)